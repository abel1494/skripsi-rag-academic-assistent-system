from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from groq import Groq
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from typing import List, Optional
import os
import json
import requests
import io
from supabase import create_client
from dotenv import load_dotenv
from datetime import datetime, timezone

load_dotenv()

# Config
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# print("Cek URL Supabase:", SUPABASE_URL)

ML_API = "https://alsyabella-ml-rag-skripsi.hf.space"

app = FastAPI()

@app.get("/api/hello")
def hello():
    return {"status": "Backend nyala, Bella!"}

supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
groq_client = Groq(api_key=GROQ_API_KEY)
timeout=60.0 

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ML integration
def get_embedding(text):
    try:
        res = requests.post(f"{ML_API}/embed", json={"text": text}, timeout=30)
        return res.json()["embedding"]
    except:
        return [0.0] * 384

def get_similarity(text1, text2):
    try:
        res = requests.post(f"{ML_API}/similarity", json={
            "text1": text1,
            "text2": text2
        }, timeout=30)
        return res.json()["similarity"]
    except:
        return 0.0

def extract_pages_from_bytes(file_bytes: bytes, filename: str):
    ext = os.path.splitext(filename)[1].lower()
    pages_data = [] 
    file_stream = io.BytesIO(file_bytes)

    if ext == ".pdf":
        reader = PdfReader(file_stream)
        for i, page in enumerate(reader.pages):
            extracted = page.extract_text()
            if extracted:
                pages_data.append({"location": f"Halaman {i+1}", "text": extracted})

    elif ext == ".docx":
        doc = Document(file_stream)
        full_text = "\n".join([para.text for para in doc.paragraphs])
        pages_data.append({"location": "Dokumen", "text": full_text})

    elif ext == ".pptx":
        prs = Presentation(file_stream)
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text + "\n"
            if slide_text.strip():
                pages_data.append({"location": f"Slide {i+1}", "text": slide_text})

    return pages_data

def chunk_text(text, chunk_size=800):
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + 1 > chunk_size and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(word)
        else:
            current_chunk.append(word)
            current_length += len(word) + 1
            
    if current_chunk:
        chunks.append(" ".join(current_chunk))
        
    return chunks

def ingest_file(file_bytes, file_name, user_id):
    pages_data = extract_pages_from_bytes(file_bytes, file_name)

    supabase.table("documents") \
        .delete() \
        .eq("file_name", file_name) \
        .eq("user_id", user_id) \
        .execute()

    data_to_insert = []
    for page_item in pages_data:
        location = page_item["location"]
        text = page_item["text"]
        
        chunks = chunk_text(text)
        
        for chunk in chunks:
            if chunk.strip():
                embedding = get_embedding(chunk)
                content_with_location = f"[{location}] {chunk}"
                
                data_to_insert.append({
                    "user_id": user_id,
                    "file_name": file_name,
                    "content": content_with_location, 
                    "embedding": embedding
                })

    if data_to_insert:
        supabase.table("documents").insert(data_to_insert).execute()

    return True


# Schema
class QuestionRequest(BaseModel):
    question: str
    file_name: Optional[List[str]] = None  
    user_id: str
    session_id: str

class QuizRequest(BaseModel):
    file_name: Optional[List[str]] = None 
    user_id: str
    num_questions: int = 5
    quiz_type: str

class AnswerCheckRequest(BaseModel):
    question_id: int
    user_answer: str
    quiz: list
    user_id: str
    file_name: str

class QuizHistoryRequest(BaseModel):
    user_id: str
    session_id: str
    quiz_type: str
    num_questions: int
    score: int
    review_data: list = []

class RegisterRequest(BaseModel):
    email: str
    password: str

class LoginRequest(BaseModel):
    email: str
    password: str


# Auth 
@app.post("/register")
def register(request: RegisterRequest):
    supabase.auth.sign_up({
        "email": request.email,
        "password": request.password
    })
    return {"message": "Register berhasil"}

@app.post("/login")
def login(request: LoginRequest):
    res = supabase.auth.sign_in_with_password({
        "email": request.email,
        "password": request.password
    })

    if res.user is None:
        return {"error": "Login gagal"}

    return {
        "message": "Login berhasil",
        "user_id": res.user.id
    }

# New Session
@app.post("/create-session")
def create_session(user_id: str):
    return supabase.table("sessions").insert({
        "user_id": user_id,
        "title": "New Chat"
    }).execute().data

# Semua Session
@app.get("/sessions")
def get_sessions(user_id: str):
    return supabase.table("sessions") \
        .select("*") \
        .eq("user_id", user_id) \
        .order("created_at", desc=True) \
        .execute().data

# Upload
@app.post("/upload")
async def upload_files(
    user_id: str = Form(...),
    files: List[UploadFile] = File(...)
):
    file_names = []

    for file in files:
        file_bytes = await file.read()

        supabase.table("files").insert({
            "user_id": user_id,
            "file_name": file.filename
        }).execute()

        ingest_file(file_bytes, file.filename, user_id)
        file_names.append(file.filename)

    return {
        "message": "Upload berhasil",
        "files": file_names
    }

# Chat
@app.post("/chat")
def chat(request: QuestionRequest):
    supabase.table("messages").insert({
        "session_id": request.session_id,
        "role": "user",
        "content": request.question
    }).execute()

    query_embedding = get_embedding(request.question)
    res = supabase.rpc("match_documents", {
        "query_embedding": query_embedding,
        "match_count": 10, 
        "user_id_param": request.user_id,
        "file_names_param": request.file_name
    }).execute()

    # sitasi
    context_texts = []
    if res.data:
        for doc in res.data:
            teks_sumber = f"--- SUMBER DARI FILE: {doc['file_name']} ---\n{doc['content']}"
            context_texts.append(teks_sumber)
    context = "\n\n".join(context_texts)

    # auto title
    existing = supabase.table("messages") \
        .select("*") \
        .eq("session_id", request.session_id) \
        .execute()

    if len(existing.data) == 1:
        title_prompt = (
            f"Berikan judul singkat maksimal 3-4 kata untuk topik pertanyaan ini: '{request.question}'. "
            "WAJIB: Berikan HANYA judulnya saja, tanpa tanda kutip, tanpa kalimat pengantar, tanpa titik di akhir."
        )
        
        title_res = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": title_prompt}],
            model="llama-3.1-8b-instant",
            temperature=0.3,
        )
        
        clean_title = title_res.choices[0].message.content.strip().replace('"', '')
        
        supabase.table("sessions").update({
            "title": clean_title
        }).eq("id", request.session_id).execute()

    messages = [
        {
            "role": "system",
            "content": (
                "Kamu adalah asisten akademik cerdas. Kamu akan diberikan potongan teks dari dokumen. "
                "Setiap potongan dokumen memiliki label 'SUMBER DARI FILE: [Nama File]' dan di dalam teksnya memuat posisi halamannya seperti '[Halaman X]'. "
                "Tugasmu: Jawab pertanyaan user BERDASARKAN konteks tersebut secara lengkap dan terstruktur. "
                "ATURAN SANGAT PENTING: Kamu WAJIB menyertakan sitasi (kutipan sumber beserta halamannya) di akhir setiap kalimat atau paragraf informasi yang kamu berikan. "
                "Gunakan format seperti ini: (Sumber: [Nama File], [Halaman X]). "
                "Jika membandingkan dua file atau lebih, sebutkan secara spesifik apa yang ada di file pertama dan apa yang ada di file kedua beserta sitasinya. "
                "Jika jawabannya tidak ada di dokumen, bilang tidak tahu. JANGAN KELUAR DARI DOKUMEN pokoknya harus yang ada di konteks yang di pilih aja. "
                "jangan kasih rekomendasi lain, pokoknya kalo ga ada di dokumen jangan di jawab dan memberikan alternatif lain untuk pencarian di web di luar itu menjadi batasan mu.\n\n"
                f"KONTEKS DOKUMEN:\n{context}"
            )
        },
        {"role": "user", "content": request.question}
    ]

    response = groq_client.chat.completions.create(
        messages=messages,
        model="llama-3.1-8b-instant",
    )

    answer = response.choices[0].message.content

    # Simpan jawaban AI
    supabase.table("messages").insert({
        "session_id": request.session_id,
        "role": "ai",
        "content": answer
    }).execute()

    return {"answer": answer}

# File 
@app.get("/files")
def list_files(user_id: str):
    return supabase.table("files") \
        .select("*") \
        .eq("user_id", user_id) \
        .execute().data

# Quiz
@app.post("/generate-quiz")
def generate_quiz(request: QuizRequest):
    try:
        print(f"--- MEMULAI GENERATE KUIS DARI {len(request.file_name)} FILE ---")
        
        all_contexts = []
        
        for fname in request.file_name:
            query_embedding = get_embedding(f"intisari materi penting dari dokumen {fname}")
            
            res = supabase.rpc("match_documents", {
                "query_embedding": query_embedding,
                "match_count": 8, 
                "user_id_param": request.user_id,
                "file_names_param": [fname] 
            }).execute()

            if res.data:
                file_context = "\n".join([d["content"] for d in res.data])
                all_contexts.append(f"--- MATERI DARI {fname} ---\n{file_context}")

        if not all_contexts:
            print("ERROR: Konteks tidak ditemukan di semua file")
            raise HTTPException(status_code=404, detail="Materi dokumen tidak ditemukan")

        context = "\n\n".join(all_contexts)[:8000]

        prompt = f"""
Anda adalah mesin pembuat kuis JSON. Buat persis {request.num_questions} soal {request.quiz_type} dari GABUNGAN MATERI di bawah ini.
PASTIKAN soal diambil secara merata dari semua dokumen yang ada di materi.
TIDAK BOLEH ADA TEKS LAIN SELAIN JSON MURNI.

Format WAJIB untuk Essay:
{{"quiz": [{{"id": 1, "question": "Apa itu X?", "answer": "X adalah..."}}]}}

Format WAJIB untuk Pilihan Ganda (PG):
{{"quiz": [{{"id": 1, "question": "Apa itu X?", "options": ["A. Y", "B. X", "C. Z", "D. W"], "answer": "B. X"}}]}}

Materi Gabungan:
{context}
        """

        response = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.1-8b-instant",
            response_format={"type": "json_object"}
        )

        raw_content = response.choices[0].message.content
        data_json = json.loads(raw_content)
        quiz_list = data_json.get("quiz", [])

        for i, q in enumerate(quiz_list):
            q["id"] = i + 1

        return {"quiz": quiz_list}

    except Exception as e:
        print(f"!!! ERROR GENERATE QUIZ: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Gagal: {str(e)}")

@app.post("/check-answer")
def check_answer(request: AnswerCheckRequest):
    q = next((x for x in request.quiz if x["id"] == request.question_id), None)
    if not q:
        return {"error": "Soal tidak ditemukan"}
        
    correct_answer = str(q["answer"]).strip()
    user_ans = str(request.user_answer).strip()

    # Cek PG
    if "options" in q:
        if user_ans.lower() == correct_answer.lower():
            return {
                "similarity": 100,
                "feedback": "Jawaban Tepat!",
                "reference_answer": correct_answer
            }
        else:
            return {
                "similarity": 0,
                "feedback": "Jawaban Salah.",
                "reference_answer": correct_answer
            }
    
    # Cek Essay
    sim = get_similarity(user_ans, correct_answer)
    percent = round(sim * 100, 2)

    feedback = "Tidak sesuai"
    if percent > 85: feedback = "Sangat sesuai"
    elif percent > 70: feedback = "Cukup sesuai"
    elif percent > 50: feedback = "Kurang tepat"

    return {
        "similarity": percent,
        "feedback": feedback,
        "reference_answer": correct_answer
    }

# Delete File
@app.delete("/delete-file")
def delete_file(file_name: str, user_id: str):
    supabase.table("files") \
        .delete() \
        .eq("file_name", file_name) \
        .eq("user_id", user_id) \
        .execute()

    supabase.table("documents") \
        .delete() \
        .eq("file_name", file_name) \
        .eq("user_id", user_id) \
        .execute()

    return {"message": "File dihapus"}

# Chat History
@app.get("/chat-history")
def get_chat_history(session_id: str):
    return supabase.table("messages") \
        .select("*") \
        .eq("session_id", session_id) \
        .order("created_at") \
        .execute().data

# Quiz History
@app.get("/quiz-history")
def get_quiz_history(session_id: str):
    try:
        res = supabase.table("quiz_history") \
            .select("*") \
            .eq("session_id", session_id) \
            .order("created_at", desc=True) \
            .execute()
        return res.data
    except Exception as e:
        print("ERROR GET HISTORY:", e)
        return []

@app.post("/save-quiz-history")
def save_quiz_history(request: QuizHistoryRequest):
    try:
        data = {
            "user_id": request.user_id,
            "session_id": request.session_id,
            "quiz_type": request.quiz_type,
            "num_questions": request.num_questions,
            "score": request.score,
            "review_data": request.review_data
        }
        res = supabase.table("quiz_history").insert(data).execute()
        return {"status": "success", "data": res.data}
    except Exception as e:
        print("ERROR SAVE HISTORY:", e)
        return {"error": str(e)}

@app.post("/update-session-time")
def update_session_time(session_id: str):
    try:
        now = datetime.now(timezone.utc).isoformat()
        res = supabase.table("sessions").update({"updated_at": now}).eq("id", session_id).execute()
        return {"status": "success"}
    except Exception as e:
        print("ERROR UPDATE WAKTU:", e)
        return {"error": str(e)}
